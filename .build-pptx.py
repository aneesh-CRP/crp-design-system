#!/usr/bin/env python3
"""Build fully editable PowerPoint decks from the rendered HTML.

Every painted box becomes a real autoshape and every text run a real text box,
positioned at the pixel coordinates Chrome laid out. 1280x720 CSS px maps 1:1
onto a 13.333 x 7.5in slide at 96dpi, so the PPTX matches the PDF closely while
staying editable — no slide images anywhere.

Run after `.pptx-*.json` have been regenerated from the browser.
"""
import json
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

DS = pathlib.Path(__file__).resolve().parent
EMU_PER_PX = 9525  # 96 dpi
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def rgbof(c):
    return RGBColor.from_string(c["hex"].upper())


def set_alpha(fill, alpha):
    """Apply real transparency.

    python-pptx has no alpha API on FillFormat — assigning `.transparency`
    silently does nothing, which turned the cover's translucent chips opaque
    white and made their white numerals disappear. Write <a:alpha> ourselves.
    """
    if alpha >= 0.99:
        return
    from pptx.oxml.ns import qn
    srgb = fill._xPr.find(qn("a:solidFill"))
    if srgb is None:
        return
    clr = srgb.find(qn("a:srgbClr"))
    if clr is None:
        return
    node = clr.makeelement(qn("a:alpha"), {"val": str(int(round(alpha * 100000)))})
    clr.append(node)


# Google Slides can only use fonts from the Google Fonts library — there is no
# upload — so Area and Market Pro always substitute there, and Slides falls back
# to Arial, which is wider and re-wraps every heading. tokens.css already names
# the sanctioned fallbacks, and both are on Google Fonts.
SLIDES_SANS, SLIDES_SCRIPT = "IBM Plex Sans", "Caveat"
SLIDES_MODE = False


def font_for(weight, family):
    """Map CSS weight onto the Area family names PowerPoint actually sees.

    The OTF name tables register ExtraBold and SemiBold as SEPARATE FAMILIES
    ("Area Normal ExtraBold", "Area Normal SemiBold"), each with a Regular
    style — only 700 is a Bold style of "Area Normal". Guessing "Area Bold"
    silently substitutes and the metrics shift enough to re-wrap headings.
    Returns (family_name, use_bold_flag).
    """
    if "Market Pro" in (family or ""):
        return (SLIDES_SCRIPT if SLIDES_MODE else "Market Pro"), False
    if "Plex Mono" in (family or ""):
        return "IBM Plex Mono", weight >= 700
    if SLIDES_MODE:
        return SLIDES_SANS, weight >= 600
    if weight >= 800:
        return "Area Normal ExtraBold", False
    if weight >= 700:
        return "Area Normal", True
    if weight >= 600:
        return "Area Normal SemiBold", False
    return "Area Normal", False


CACHE = DS / ".pptx-imgcache"


def downscale(src, box_w, box_h, factor=2):
    """Embed images at ~2x their placed size.

    Sponsor marks are rasterised at 600px tall but placed at ~30-56px, so the
    file carried roughly an order of magnitude more pixels than it renders.
    Caps each image at `factor` x its box and caches the result.
    """
    try:
        from PIL import Image
        im = Image.open(src)
        target_w, target_h = int(box_w * factor), int(box_h * factor)
        if im.width <= target_w * 1.15 and im.height <= target_h * 1.15:
            return src
        scale = max(target_w / im.width, target_h / im.height)
        new = (max(1, int(im.width * scale)), max(1, int(im.height * scale)))
        CACHE.mkdir(exist_ok=True)
        out = CACHE / f"{src.stem}-{new[0]}x{new[1]}{src.suffix}"
        if not out.exists():
            im.convert("RGBA" if src.suffix.lower() == ".png" else "RGB") \
              .resize(new, Image.LANCZOS).save(out, optimize=True)
        return out
    except Exception:
        return src


def script_image(text, size_px, color_hex):
    """Rasterise the Market Pro tagline for the Google Slides build.

    Market Pro is a compact script; at the same point size Caveat or Arial set
    two to three times wider, so the cover lockup overflows and collides. It is
    a brand mark rather than editable copy, so the Slides variant ships it as a
    picture and keeps the layout exact.
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        return None
    face = pathlib.Path.home() / "Library/Fonts/MarketPro.otf"
    if not face.exists():
        face = DS / "assets/fonts/MarketPro.otf"
    if not face.exists():
        return None
    S = 3  # supersample
    try:
        font = ImageFont.truetype(str(face), int(size_px * S))
    except Exception:
        return None
    probe = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    box = probe.textbbox((0, 0), text, font=font)
    w, h = max(1, box[2] - box[0]), max(1, box[3] - box[1])
    pad = int(size_px * S * 0.2)
    img = Image.new("RGBA", (w + pad * 2, h + pad * 2), (0, 0, 0, 0))
    ImageDraw.Draw(img).text((pad - box[0], pad - box[1]), text, font=font,
                             fill="#" + color_hex.upper())
    CACHE.mkdir(exist_ok=True)
    safe = "".join(c for c in text if c.isalnum())[:24]
    out = CACHE / f"script-{safe}-{int(size_px)}-{color_hex}.png"
    img.resize((img.width // S, img.height // S), Image.LANCZOS).save(out)
    return out, img.width // S, img.height // S


def add_shape(slide, s):
    rounded = s["radius"] >= 6
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        px(s["x"]), px(s["y"]), px(s["w"]), px(s["h"]))
    if rounded:
        try:  # adjustment is a fraction of the shorter side
            shp.adjustments[0] = min(0.5, s["radius"] / max(1.0, min(s["w"], s["h"])))
        except (IndexError, ValueError):
            pass
    # Cover panels are CSS linear-gradients. backgroundColor reports transparent
    # for those, so without this the navy panel vanishes and its white text
    # lands invisibly on a white slide.
    if s.get("grad"):
        stops = s["grad"]
        shp.fill.gradient()
        gs = shp.fill.gradient_stops
        gs[0].color.rgb = RGBColor.from_string(stops[0].upper())
        gs[0].position = 0.0
        gs[1].color.rgb = RGBColor.from_string(stops[-1].upper())
        gs[1].position = 1.0
        try:
            shp.fill.gradient_angle = float(s.get("gradAngle", 135))
        except (ValueError, TypeError):
            pass
    elif s["fill"]:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgbof(s["fill"])
        set_alpha(shp.fill, s["fill"]["a"])
    else:
        shp.fill.background()
    # A CSS box can have four different borders; PowerPoint has one outline.
    # Use the heaviest edge, and lay a separate bar for a lone accent edge.
    widths = s["bw"]
    heaviest = max(range(4), key=lambda i: widths[i])
    uniform = len({round(w, 1) for w in widths}) == 1 and widths[0] > 0.4
    if uniform and s["bc"][0]:
        shp.line.color.rgb = rgbof(s["bc"][0])
        shp.line.width = Pt(widths[0] * 0.75)
    elif widths[heaviest] > 0.4 and s["bc"][heaviest]:
        shp.line.fill.background()
        edge = ["top", "right", "bottom", "left"][heaviest]
        w, h = s["w"], s["h"]
        geom = {"top": (s["x"], s["y"], w, widths[0]),
                "bottom": (s["x"], s["y"] + h - widths[2], w, widths[2]),
                "left": (s["x"], s["y"], widths[3], h),
                "right": (s["x"] + w - widths[1], s["y"], widths[1], h)}[edge]
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(geom[0]), px(geom[1]),
                                     px(max(geom[2], 0.8)), px(max(geom[3], 0.8)))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgbof(s["bc"][heaviest])
        bar.line.fill.background()
        bar.shadow.inherit = False
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, t):
    if SLIDES_MODE and "Market Pro" in (t.get("family") or ""):
        txt = "".join(r["t"] for r in t["runs"]).strip()
        col = next((r["color"]["hex"] for r in t["runs"] if r.get("color")), "FF9933")
        made = script_image(txt, t["size"], col)
        if made:
            path, iw, ih = made
            scale = min(t["w"] / iw, t["h"] / ih) if iw and ih else 1
            w, h = iw * scale, ih * scale
            return slide.shapes.add_picture(
                str(path), px(t["x"] + (t["w"] - w) / 2), px(t["y"] + (t["h"] - h) / 2),
                px(w), px(h))

    # Give the box slack so PowerPoint's slightly different metrics do not
    # re-wrap a line that fits in the browser. Centred text grows symmetrically.
    pad = 26 if t["align"] == "center" else 14
    dx = pad / 2 if t["align"] == "center" else 2
    box = slide.shapes.add_textbox(px(t["x"] - dx), px(t["y"] - 3),
                                   px(t["w"] + pad), px(t["h"] + 8))
    tf = box.text_frame
    tf.word_wrap = t["h"] > t["lh"] * 1.45
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para = tf.paragraphs[0]
    para.alignment = ALIGN.get(t["align"], PP_ALIGN.LEFT)
    para.line_spacing = Pt(t["lh"] * 0.75)
    runs = list(t["runs"])
    marker = (t.get("before") or "").strip()
    if marker and runs:
        lead = dict(runs[0])
        lead["t"] = marker + " "
        if t.get("beforeColor"):
            lead["color"] = t["beforeColor"]
        runs.insert(0, lead)
    for r in runs:
        text = r["t"]
        if not text:
            continue
        if text == "\n":
            para = tf.add_paragraph()
            para.alignment = ALIGN.get(t["align"], PP_ALIGN.LEFT)
            para.line_spacing = Pt(t["lh"] * 0.75)
            continue
        run = para.add_run()
        run.text = text.upper() if t["upper"] else text
        href = r.get("href") or t.get("href")
        if href:
            run.hyperlink.address = href
        f = run.font
        fam, bold = font_for(r["weight"], t.get("family"))
        f.name = fam
        f.size = Pt(r["size"] * 0.75)
        f.bold = bold
        f.italic = r["italic"]
        if r["color"]:
            f.color.rgb = rgbof(r["color"])
    return box


def add_image(slide, im):
    if "quat" in (im.get("cls") or "") or "stencil" in im["src"]:
        return None  # decorative watermark; rasterises to an opaque blob
    # Forcing lazy images to load rewrites src to an absolute file:// URL, so
    # accept both that and the authored relative path.
    raw = im["src"]
    if raw.startswith("file://"):
        from urllib.parse import unquote, urlparse
        src = pathlib.Path(unquote(urlparse(raw).path))
    else:
        src = (DS / raw.lstrip("./")).resolve()
    if not src.exists():
        return None
    if src.suffix.lower() == ".svg":  # PowerPoint cannot place SVG reliably
        png = src.with_suffix(".png")
        if not png.exists():
            return None
        src = png
    # PowerPoint has no CSS filters, so the sponsor wall's grayscale(100%)
    # has to come from a pre-desaturated file.
    if im.get("gray"):
        grey = src.with_name(src.stem + "-gray.png")
        if grey.exists():
            src = grey
    src = downscale(src, im["w"], im["h"])
    try:
        pic = slide.shapes.add_picture(str(src), px(im["x"]), px(im["y"]),
                                       px(im["w"]), px(im["h"]))
    except Exception:
        return None
    # object-fit: cover crops the overflow; add_picture would stretch instead.
    if im.get("fit") == "cover":
        try:
            from PIL import Image
            nw, nh = Image.open(src).size
            box_ar, img_ar = im["w"] / im["h"], nw / nh
            if img_ar > box_ar:                       # image too wide: crop sides
                keep = box_ar / img_ar
                pic.crop_left = pic.crop_right = (1 - keep) / 2
            elif img_ar < box_ar:                     # too tall: crop top/bottom
                keep = img_ar / box_ar
                pic.crop_top = pic.crop_bottom = (1 - keep) / 2
        except Exception:
            pass
    return pic


def build(json_name, out_name):
    data = json.loads((DS / json_name).read_text())
    prs = Presentation()
    prs.slide_width, prs.slide_height = px(1280), px(720)
    blank = prs.slide_layouts[6]
    for sl in data["slides"]:
        s = prs.slides.add_slide(blank)
        for shp in sl["shapes"]:
            add_shape(s, shp)
        for im in sl["images"]:
            add_image(s, im)
        for t in sl["texts"]:
            add_text(s, t)
    out = pathlib.Path.home() / out_name
    prs.save(str(out))
    return out, len(data["slides"])


if __name__ == "__main__":
    import sys
    decks = [(".pptx-overview.json", "CRP-Executive-Overview-Aug-2026"),
             (".pptx-derm.json", "CRP-Dermatology-Site-Capabilities"),
             (".pptx-alz.json", "CRP-Alzheimers-Site-Capabilities")]
    for src, stem in decks:
        path, n = build(src, stem + ".pptx")
        print(f"  {path.name:<50} {n} slides")
    if "--slides" in sys.argv:
        SLIDES_MODE = True
        globals()["SLIDES_MODE"] = True
        for src, stem in decks:
            path, n = build(src, stem + "-GoogleSlides.pptx")
            print(f"  {path.name:<50} {n} slides  (Google Fonts)")
