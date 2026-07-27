#!/usr/bin/env python3
"""Layout helpers for a text-native ICON deck.

Everything here emits real PowerPoint shapes and text frames — no page images —
so the deck can be edited in PowerPoint or Keynote. Fonts fall back to Arial
because Area Normal is licensed and will not be on a recipient's machine.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x07, 0x20, 0x61)
BLUE   = RGBColor(0x18, 0x43, 0xAD)
ORANGE = RGBColor(0xFF, 0x99, 0x33)
CYAN   = RGBColor(0xA2, 0xDC, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x5E, 0x76, 0x94)
RULE   = RGBColor(0xD9, 0xE2, 0xEF)
TRACK  = RGBColor(0xE5, 0xEC, 0xF5)
BG     = RGBColor(0xF4, 0xF8, 0xFF)

FONT = "Arial"
W, H = Inches(13.333), Inches(7.5)
M    = Inches(0.62)                     # page margin
CONTENT_W = W - 2 * M


def _tf(shape, text, size, *, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        space=0, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = color
    if space:
        p.space_after = Pt(space)
    return tf


def textbox(slide, x, y, w, h, text, size, **kw):
    box = slide.shapes.add_textbox(x, y, w, h)
    _tf(box, text, size, **kw)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        shp.adjustments[0] = 0.06
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def head(slide, eyebrow, title, lede=None):
    """Eyebrow, action title, optional standfirst — centred, like the HTML."""
    y = Inches(0.42)
    if eyebrow:
        textbox(slide, M, y, CONTENT_W, Inches(0.25),
                eyebrow.upper(), 11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        y += Inches(0.34)
    tb = slide.shapes.add_textbox(M, y, CONTENT_W, Inches(0.9))
    _tf(tb, title, 27, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    y += Inches(0.62) + (Inches(0.42) if len(title) > 62 else Emu(0))
    if lede:
        textbox(slide, M + Inches(1.1), y, CONTENT_W - Inches(2.2), Inches(0.5),
                lede, 13, color=GREY, align=PP_ALIGN.CENTER)
        y += Inches(0.46)
    return y + Inches(0.22)


def footer(slide, left, right):
    textbox(slide, M, H - Inches(0.52), Inches(7), Inches(0.24),
            left.upper(), 8, color=GREY, space=0)
    textbox(slide, W - M - Inches(2.6), H - Inches(0.52), Inches(2.6), Inches(0.24),
            right.upper(), 8, color=GREY, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, accent=BLUE):
    """White card with the house accent rule across the top."""
    body = rect(slide, x, y, w, h, fill=WHITE, line=RULE)
    bar = rect(slide, x, y, w, Pt(3), fill=accent, radius=False)
    bar.line.fill.background()
    return body


def stat(slide, x, y, w, value, label, *, size=40, color=NAVY, label_size=11):
    textbox(slide, x, y, w, Inches(0.6), value, size, bold=True, color=color)
    textbox(slide, x, y + Inches(0.62), w, Inches(0.7), label, label_size, color=NAVY)


def bar(slide, x, y, w, pct, *, color=ORANGE, height=Pt(7), target=None):
    track = rect(slide, x, y, w, height, fill=TRACK, radius=False)
    track.line.fill.background()
    if pct > 0:
        fill = rect(slide, x, y, Emu(int(w * pct / 100)), height, fill=color, radius=False)
        fill.line.fill.background()
    if target is not None:
        t = rect(slide, x + Emu(int(w * target / 100)), y - Pt(2.5), Pt(1.5),
                 height + Pt(5), fill=NAVY, radius=False)
        t.line.fill.background()


def bullets(slide, x, y, w, items, *, size=11.5, gap=0.30, dot=BLUE):
    """Bulleted lines with a drawn dot, so the text stays a plain editable run."""
    cy = y
    for it in items:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy + Inches(0.055), Pt(4), Pt(4))
        d.fill.solid(); d.fill.fore_color.rgb = dot
        d.line.fill.background(); d.shadow.inherit = False
        tb = slide.shapes.add_textbox(x + Inches(0.16), cy, w - Inches(0.16), Inches(0.3))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        for seg, bold in it if isinstance(it, list) else [(it, False)]:
            r = p.add_run(); r.text = seg
            r.font.name, r.font.size, r.font.bold = FONT, Pt(size), bold
            r.font.color.rgb = NAVY
        cy += Inches(gap) + Inches(0.14) * (len(str(it)) // 78)
    return cy


def section_cover(slide, num, name, blurb, value, label):
    bg = rect(slide, 0, 0, W, H, fill=NAVY, radius=False)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    x = Inches(0.95)
    textbox(slide, x, Inches(1.55), Inches(6), Inches(0.3),
            "SECTION %02d" % num, 11, bold=True, color=CYAN)
    textbox(slide, x, Inches(1.95), Inches(8), Inches(0.9), name, 40, bold=True, color=WHITE)
    ln = rect(slide, x, Inches(2.95), Inches(6.4), Pt(0.75), fill=RGBColor(0x3E,0x5C,0xA8), radius=False)
    ln.line.fill.background()
    textbox(slide, x, Inches(3.18), Inches(6.3), Inches(0.9), blurb, 13.5, color=RGBColor(0xD5,0xDE,0xF2))
    textbox(slide, x, Inches(4.35), Inches(6), Inches(0.9), value, 46, bold=True, color=ORANGE)
    textbox(slide, x, Inches(5.15), Inches(6), Inches(0.4), label, 11.5, color=RGBColor(0xC3,0xD0,0xEA))
    textbox(slide, x, H - Inches(0.52), Inches(6), Inches(0.24),
            name.upper(), 8, color=RGBColor(0x8E,0xA4,0xCE))
