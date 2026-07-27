#!/usr/bin/env python3
"""Build the ICON deck as an editable PowerPoint.

Every slide is real shapes and text — no page images — so copy can be changed
in PowerPoint or Keynote. Content lives in DECK below; the layout functions
above it are the whole design vocabulary.
"""
import pathlib, sys
sys.path.insert(0, str(pathlib.Path(__file__).parent))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import pptx_kit as K
from pptx_kit import (NAVY, BLUE, ORANGE, WHITE, GREY, RULE, W, H, M, CONTENT_W,
                      textbox, rect, card, head, footer, stat, bar, bullets,
                      section_cover)

HOME = pathlib.Path.home()
DS = HOME / "crp-design-system"
OUT = HOME / "CRP-ICON-Elite-Sites-Metabolic-Neurology.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
PAGE = [0]

def new(note=""):
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK)
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s

def pageno():
    return "%02d" % PAGE[0]

# ── layouts ────────────────────────────────────────────────────────────────

def slide_stats(eyebrow, title, lede, items, note, *, size=34):
    s = new(title)
    y = head(s, eyebrow, title, lede)
    n = len(items)
    cw = (CONTENT_W - Inches(0.25) * (n - 1)) / n
    for i, (v, k) in enumerate(items):
        stat(s, M + i * (cw + Inches(0.25)), y + Inches(0.35), cw, v, k,
             size=size, color=ORANGE if i == 0 else NAVY, label_size=11)
    if note:
        textbox(s, M, H - Inches(1.35), CONTENT_W, Inches(0.5), note, 10.5, color=GREY)
    footer(s, eyebrow, pageno())
    return s

def slide_table(eyebrow, title, lede, cols, rows, note, widths):
    s = new(title)
    y = head(s, eyebrow, title, lede)
    xs, x = [], M
    for wd in widths:
        xs.append(x); x += Inches(wd)
    for cx, c in zip(xs, cols):
        if c:
            textbox(s, cx, y, Inches(3), Inches(0.24), c.upper(), 9, bold=True, color=BLUE)
    y += Inches(0.32)
    ln = rect(s, M, y, CONTENT_W, Pt(0.75), fill=RULE, radius=False); ln.line.fill.background()
    y += Inches(0.12)
    for r in rows:
        for i, (cx, cell) in enumerate(zip(xs, r)):
            if cell is None:
                continue
            if isinstance(cell, tuple):                    # ("bar", pct, target)
                bar(s, cx, y + Inches(0.1), Inches(widths[i] - 0.25), cell[1],
                    target=cell[2], color=cell[3] if len(cell) > 3 else ORANGE)
            else:
                textbox(s, cx, y, Inches(widths[i] - 0.15), Inches(0.3), str(cell),
                        11.5, bold=(i == 0), color=NAVY if i == 0 else GREY)
        y += Inches(0.46)
    if note:
        textbox(s, M, y + Inches(0.18), CONTENT_W, Inches(0.5), note, 10.5, color=GREY)
    footer(s, eyebrow, pageno())
    return s

def slide_panels(eyebrow, title, lede, left, right, note=None):
    s = new(title)
    y = head(s, eyebrow, title, lede)
    pw = (CONTENT_W - Inches(0.28)) / 2
    ph = H - y - Inches(1.15)
    for idx, (hdr, items, dark) in enumerate((left, right)):
        x = M + idx * (pw + Inches(0.28))
        if dark:
            rect(s, x, y, pw, ph, fill=NAVY)
        else:
            card(s, x, y, pw, ph)
        textbox(s, x + Inches(0.24), y + Inches(0.22), pw - Inches(0.48), Inches(0.25),
                hdr.upper(), 9, bold=True, color=K.CYAN if dark else BLUE)
        cy = y + Inches(0.58)
        for it in items:
            tb = textbox(s, x + Inches(0.42), cy, pw - Inches(0.68), Inches(0.4), it, 11.5,
                         color=WHITE if dark else NAVY)
            d = s.shapes.add_shape(K.MSO_SHAPE.OVAL, x + Inches(0.24), cy + Inches(0.07), Pt(4), Pt(4))
            d.fill.solid(); d.fill.fore_color.rgb = ORANGE if dark else BLUE
            d.line.fill.background(); d.shadow.inherit = False
            cy += Inches(0.30) + Inches(0.17) * (len(it) // 62)
    if note:
        textbox(s, M, H - Inches(1.0), CONTENT_W, Inches(0.4), note, 10.5, color=GREY)
    footer(s, eyebrow, pageno())
    return s

def slide_list(eyebrow, title, lede, left_items, right_title, right_rows, note=None):
    s = new(title)
    y = head(s, eyebrow, title, lede)
    lw = CONTENT_W * 0.42
    rw = CONTENT_W - lw - Inches(0.3)
    ph = H - y - Inches(1.15)
    card(s, M, y, lw, ph)
    bullets(s, M + Inches(0.24), y + Inches(0.28), lw - Inches(0.48), left_items, gap=0.36)
    x2 = M + lw + Inches(0.3)
    card(s, x2, y, rw, ph)
    textbox(s, x2 + Inches(0.24), y + Inches(0.22), rw - Inches(0.48), Inches(0.25),
            right_title.upper(), 9, bold=True, color=BLUE)
    cy = y + Inches(0.6)
    for r in right_rows:
        textbox(s, x2 + Inches(0.24), cy, rw - Inches(0.48), Inches(0.28), r, 10.5)
        cy += Inches(0.31)
    if note:
        textbox(s, M, H - Inches(1.0), CONTENT_W, Inches(0.4), note, 10.5, color=GREY)
    footer(s, eyebrow, pageno())
    return s

def slide_people(eyebrow, title, lede, people, note, cols=4):
    s = new(title)
    y = head(s, eyebrow, title, lede)
    cw = (CONTENT_W - Inches(0.2) * (cols - 1)) / cols
    for i, (name, role, spec) in enumerate(people):
        x = M + (i % cols) * (cw + Inches(0.2))
        cy = y + (i // cols) * Inches(1.15)
        card(s, x, cy, cw, Inches(0.98), accent=BLUE if i % 2 == 0 else ORANGE)
        textbox(s, x + Inches(0.18), cy + Inches(0.18), cw - Inches(0.36), Inches(0.26),
                name, 11.5, bold=True)
        textbox(s, x + Inches(0.18), cy + Inches(0.45), cw - Inches(0.36), Inches(0.22),
                role, 9.5, color=GREY)
        if spec:
            textbox(s, x + Inches(0.18), cy + Inches(0.66), cw - Inches(0.36), Inches(0.22),
                    spec.upper(), 8.5, bold=True, color=BLUE)
    if note:
        textbox(s, M, H - Inches(1.05), CONTENT_W, Inches(0.5), note, 10.5, color=GREY)
    footer(s, eyebrow, pageno())
    return s

# ── deck ───────────────────────────────────────────────────────────────────

SECTIONS = [
    ("Cardiometabolic", "Our anchor area, and where we have delivered", "5 of 5", "recent targets met or beaten"),
    ("Neurology", "Migraine and MCI to mild Alzheimer's", "4", "certified raters on MMSE, CDR, ADAS-Cog, FAQ"),
    ("Derm & Rheum", "The adjacent bench, on the same model", "9", "recent trials, 2024 to 2025"),
    ("How we deliver", "Why the results repeat: recruitment, speed, quality", "<21", "days from selection to activation"),
    ("Track record", "Thirty years, the bench, and the sponsors who return", "15", "protocols placed by Lilly alone"),
    ("Partnership", "What we ask, and what we commit to", "1", "contract, one team, two addresses"),
]

# 1 cover
s = new("Cover")
bg = rect(s, 0, 0, W, H, fill=NAVY, radius=False); bg.line.fill.background()
textbox(s, Inches(0.95), Inches(1.8), Inches(6), Inches(0.3),
        "CLINICAL RESEARCH PHILADELPHIA", 11.5, bold=True, color=K.CYAN)
textbox(s, Inches(0.95), Inches(2.25), Inches(7.4), Inches(1.9),
        "Metabolic and\nNeurology experience", 40, bold=True, color=WHITE)
textbox(s, Inches(0.95), Inches(4.15), Inches(6.4), Inches(0.4),
        "Two sites in Greater Philadelphia, run as one operation.", 14, color=K.RGBColor(0xD5,0xDE,0xF2))
textbox(s, Inches(0.95), Inches(4.95), Inches(4), Inches(0.3), "PREPARED FOR", 10.5, bold=True, color=K.CYAN)
textbox(s, Inches(0.95), Inches(5.28), Inches(6), Inches(0.35), "Derek Griggs, BSN, RN", 15, bold=True, color=WHITE)
textbox(s, Inches(0.95), Inches(5.66), Inches(6), Inches(0.3),
        "ICON Elite Sites Program · July 2026", 11.5, color=K.RGBColor(0xB9,0xC7,0xE6))
_icon = DS / "assets/logos/iconplc-white.png"
if _icon.exists():
    s.shapes.add_picture(str(_icon), Inches(0.95), Inches(6.05), width=Inches(1.35))
_logo = DS / "assets/logos/crp-white.png"
if _logo.exists():
    s.shapes.add_picture(str(_logo), W - Inches(4.9), Inches(3.1), width=Inches(3.6))
textbox(s, Inches(0.95), H - Inches(0.52), Inches(4), Inches(0.24),
        "PHILLYRESEARCH.COM", 8, color=K.RGBColor(0x8E,0xA4,0xCE))

# 2 home
s = new("We beat the enrollment goal at both sites, and the data holds up.")
y = head(s, "Clinical Research Philadelphia · for ICON Elite Sites",
         "We beat the enrollment goal at both sites, and the data holds up.")
top = [("3.5×", "the enrollment goal, on our largest recent trial", ORANGE),
       ("5 of 5", "recent cardiometabolic targets met or beaten", NAVY),
       ("83%", "average retention across those trials", NAVY),
       ("0", "FDA 483s and no sponsor audit findings — our data has never been cited", NAVY)]
cw = (CONTENT_W - Inches(0.6)) / 4
for i, (v, k, c) in enumerate(top):
    stat(s, M + i * (cw + Inches(0.2)), y, cw, v, k, size=32, color=c, label_size=10)
y += Inches(1.7)
ln = rect(s, M, y, CONTENT_W, Pt(0.75), fill=RULE, radius=False); ln.line.fill.background()
y += Inches(0.24)
cw3 = (CONTENT_W - Inches(0.4)) / 3
for i, (name, desc, v, k) in enumerate(SECTIONS):
    cx = M + (i % 3) * (cw3 + Inches(0.2))
    cy = y + (i // 3) * Inches(1.22)
    card(s, cx, cy, cw3, Inches(1.06), accent=BLUE if i % 2 == 0 else ORANGE)
    textbox(s, cx + Inches(0.18), cy + Inches(0.15), cw3 - Inches(0.36), Inches(0.26), name, 12.5, bold=True)
    textbox(s, cx + Inches(0.18), cy + Inches(0.43), cw3 - Inches(0.36), Inches(0.36), desc, 9.5, color=GREY)
    textbox(s, cx + Inches(0.18), cy + Inches(0.78), cw3 - Inches(0.36), Inches(0.24),
            "%s %s" % (v, k), 9, bold=True, color=ORANGE)
footer(s, "Start anywhere", pageno())

def cover(i):
    n, d, v, k = SECTIONS[i]
    s = new(n)
    section_cover(s, i + 1, n, d, v, k)

# §1 cardiometabolic
cover(0)
slide_table("Cardiometabolic", "Five for five. Every recent target met or beaten.",
  "Our five most recent cardiometabolic trials. Bar length is delivery against the sponsor's allocation; counts are from our CTMS.",
  ["Study", "Versus target", None, "Retention", None],
  [["Lp(a) ASCVD outcomes", ("bar", 100, 28.8), "347%", ("bar", 95, None, BLUE), "95%"],
   ["ASCVD + CKD",          ("bar", 53.6, 28.8), "186%", ("bar", 85, None, BLUE), "85%"],
   ["T2DM + CV risk",       ("bar", 51.9, 28.8), "180%", ("bar", 85, None, BLUE), "85%"],
   ["T2DM — GLP first line",("bar", 40.3, 28.8), "140%", ("bar", 80, None, BLUE), "80%"],
   ["T2DM — GLP vs insulin",("bar", 28.8, 28.8), "100%", ("bar", 70, None, BLUE), "70%"]],
  "Eleven metabolic protocols in total — Lilly, AstraZeneca, Amgen, Novo Nordisk. Blue bar is retention.",
  [3.6, 3.6, 1.1, 1.6, 0.9])
slide_panels("Cardiometabolic · how the funnel converted",
  "Fifty-two randomizations cost 1,180 worked patients, and we absorbed that.",
  "Lilly J3L-MC-EZEF. 1,180 identified and worked · 262 reached a screening outcome · 210 screen failed · 52 randomized.",
  ("What made it work",
   ["Lp(a) testing in house — eligibility settled before screening, not during it.",
    "Cardiology referral base built over years, not bought per study.",
    "One coordinator queue, so every patient has a single owner.",
    "Both sites eligible, drawing from separate catchments."], False),
  ("What that funnel produced",
   ["52 randomized on a single protocol, against an allocation of 15.",
    "1,180 patients worked from our own database and referral network, with no per-study advertising spend.",
    "95% retention — the patients we randomized stayed on study."], True))
slide_list("Cardiometabolic capability",
  "Four specialists and nine recent trials mean the next protocol starts with a team that has run it.",
  "Our anchor therapeutic area — the one where we have both the investigators and the referral base already in place.",
  ["Joseph Heether, MD — Cardiology", "Taher Modarressi, MD, FNLA — Endocrinology / Lipidology",
   "Hal Ganzman, DO — Sub-I, Internal Medicine", "Vinod Rustgi, MD, MBA — Hepatology / Liver",
   "FibroScan, ultrasound, EKG, centrifuges and infusion pumps on site",
   "Lp(a) and lipid panels run in house before a patient reaches screening"],
  "Recent research experience",
  ["LDL-cholesterol reduction in high-risk ASCVD · Phase III · 2025, ongoing",
   "GIP/GLP-1 receptor agonist in MASLD · Phase III · 2025, ongoing",
   "ASCVD risk reduction with elevated lipoprotein · Phase III · 2025, ongoing",
   "Cardiovascular outcomes in chronic kidney disease · Phase III · 2025, ongoing",
   "Lipoprotein(a) lowering, primary and secondary prevention · Phase III · 2025",
   "Oral GLP-1 therapy for T2D and cardiometabolic risk · Phase III · 2024",
   "GLP-1 weight management in adults · Phase III · 2024",
   "Cardiometabolic prevention in insulin resistance · Phase II · 2024",
   "LDL-cholesterol reduction in high-risk ASCVD · Phase III · 2024"])

# §2 neurology
cover(1)
slide_stats("Neurology", "Four certified raters are why an early Alzheimer's protocol can be placed here at all.",
  "Migraine and MCI to mild Alzheimer's. We do not run psychiatry.",
  [("4", "certified raters on MMSE, CDR, ADAS-Cog and FAQ"),
   ("8", "protocols across six sponsors"),
   ("5", "separate migraine assets"),
   ("4", "open and enrolling now")],
  "Rater certification is the constraint most sites fail on. Ours are trained and current on the instruments these trials run, with the neurology, geriatrics and primary-care referral base behind them.")
slide_table("Neurology · protocols at our sites",
  "Biohaven, AbbVie and Lilly have each placed neurology work with us.",
  "Eight protocols in total. The four that carried the enrollment are below — 23 randomized from roughly 2,600 patients worked.",
  ["Protocol", "Sponsor and indication", "Randomized", "Patients worked", "Status"],
  [["BHV3000-405", "Biohaven — episodic migraine", "8", "1,304", "Closed"],
   ["M23-714", "AbbVie — menstrual migraine", "7", "486", "Enrolling"],
   ["J1G-MC-LAKI", "Lilly — Alzheimer's disease", "6", "494", "Closed"],
   ["KAL-K304-P001", "Acute migraine", "2", "284", "Enrolling"]],
  "A migraine franchise sponsors return to. Five separate assets across six sponsors, plus a Lilly Alzheimer's protocol.",
  [2.6, 4.3, 1.6, 1.9, 1.7])
slide_list("Early Alzheimer's · depth",
  "Our investigators have run six Alzheimer's trials, so the raters are not starting cold.",
  "What the raters sit on top of: the team, the referral base and the trials our investigators have run.",
  ["4 certified raters trained on MMSE, CDR, ADAS-Cog and FAQ",
   "5 clinical research coordinators experienced in cognitive assessment visits",
   "Board-certified neurologist and geriatrician reached through our provider network",
   "Referral base across neurology, geriatrics and primary care",
   "Our Alzheimer's pre-screening pathway worked 494 patients for a Lilly protocol"],
  "Investigator experience in Alzheimer's disease",
  ["Early Alzheimer's monoclonal antibody · Phase III · 2025, ongoing",
   "Sensory stimulation for Alzheimer's · Phase III · 2024",
   "Alzheimer's disease prevention · Phase III · 2022",
   "Early-stage Alzheimer's immunotherapy · Phase I/II · 2022",
   "Alzheimer's monoclonal antibody · Phase III · 2021",
   "Cognitive health and metabolic intervention · Phase III · 2021"],
  "Investigator experience, including trials run by our investigators and network-affiliated specialists.")

# §3 derm & rheum
cover(2)
slide_list("Dermatology and rheumatology",
  "Nine derm and rheum trials run on the same coordinators, so placement here needs no new set-up.",
  "The same coordinators, regulatory function and referral engine — placement here needs no new set-up.",
  ["Michael Tomeo, MD — Dermatology", "Lawrence Leventhal, MD — Rheumatology",
   "Parth Patel, PA — Dermatology", "Lolita Vaughan, CRNP — Derm / Rheum",
   "IGA and PASI scoring performed on site",
   "Biologic-experienced and JAK-experienced patients already identified in our database"],
  "Recent research experience, 2024 to 2025",
  ["Atopic dermatitis, unresponsive to biologics and JAK inhibitors · Phase III",
   "Plaque psoriasis, oral IL-23 receptor antagonist · Phase III",
   "Chronic spontaneous urticaria, antihistamine-refractory · Phase III",
   "Lichen simplex chronicus · Phase III",
   "Atopic dermatitis, JAK1 inhibitor in biologic-experienced patients · Phase III",
   "Moderate-to-severe eczema, OX40 inhibition · Phase II",
   "Rheumatoid arthritis, methotrexate-inadequate responders · Phase III",
   "Sjögren's syndrome, BTK pathway · Phase III",
   "Systemic lupus erythematosus, ROCK inhibition · Phase III"])

# §4 how we deliver
cover(3)
slide_stats("Recruitment", "191 provider partnerships put patients in front of us before recruitment starts.", None,
  [("191", "referring practices across PA and NJ"),
   ("1.5M", "patients in those provider panels"),
   ("37k", "patients in our own database, re-screened against every protocol")],
  "Referring physicians send us diagnosed patients directly, and the network does not reset when a study closes — so the second protocol in an area enrolls faster than the first.", size=44)
slide_stats("Diversity in enrollment",
  "One in three participants is Black or African American, without a separate diversity spend.",
  "Every participant with a demographics record in our EDC, 2021 to date. n = 2,010 with race reported.",
  [("32%", "Black or African American, against 13.7% of the US population"),
   ("41%", "not White, where race is reported"),
   ("9%", "Hispanic or Latino, of those reporting ethnicity")],
  "Both sites sit inside the populations sponsors struggle to reach, so representative enrollment is the default rather than a campaign.", size=40)
slide_table("Performance proof", "We met or beat the enrollment target on nine of our last ten trials.",
  "Cycle times measured against NIH guidelines (PubMed 30227522).",
  ["Days", "", None, None, None],
  [["<7", "From greenlight to first participant screened", None, None, None],
   ["3", "From start-up packet to IRB submission", None, None, None],
   ["<7", "From draft budget to contract finalized", None, None, None],
   ["3", "Average query resolution", None, None, None],
   ["1", "Average EDC data entry — CRIO source-to-EDC", None, None, None]],
  "95% of recent studies met or exceeded the sponsor's enrollment target. In our cardiometabolic trials we enrolled 340% and 186% of target.",
  [1.4, 9.0, 0.5, 0.5, 0.5])
slide_panels("What our technology gets you",
  "Every visit is captured electronically, so fewer discrepancies ever reach you.",
  "CRIO is our source of truth. What we built on top reads it continuously so problems surface while they are still fixable.",
  ("For your monitors and study managers",
   ["Source entered the day of the visit — monitoring visits review data instead of waiting on it.",
    "Enrollment, screen-fail reasons and visit compliance available on demand.",
    "Prior inclusion and exclusion results carry forward instead of re-working a patient."], False),
  ("For your risk",
   ["Missed visits and drifting screen-fail rates flag while they are still fixable.",
    "Every outbound patient message reviewed before it is sent.",
    "Per-action audit log, BAA-ready, HIPAA-grade."], True))
slide_panels("Risk profile", "Nothing in our risk profile should slow a placement decision.",
  "Zero 483s in thirty years, no findings at the last sponsor audit, standing IRBs, and capacity to start now.",
  ("Compliance and quality",
   ["30 years of continuous operation, zero FDA Form 483s.",
    "Last sponsor audit December 2024 — no findings.",
    "SOPs maintained for ICH-GCP E6 (R3).",
    "Standing IRB relationships: Advarra, WCG, Sterling."], False),
  ("Capacity and security",
   ["24 currently enrolling protocols, two sites running in parallel.",
    "In-house budget and contracting team.",
    "Encrypted, role-based access with every action logged.",
    "Direct EHR connection for medical records, BAA-ready from kickoff."], True))
slide_people("Operations and coordinators",
  "Coordinators average eight years here, so your study sees the same faces throughout.",
  "Continuity of staff across visits — your study sees the same faces from screening to closeout.",
  [("Liz Usmiani", "Recruitment & Ops Lead", ""), ("Leticia Denato, MPH", "Operations Manager", ""),
   ("Stacey Scott, RMA", "Coordinator", ""), ("Ruby Pereira", "Coordinator", ""),
   ("Cady Chilensky", "Coordinator", ""), ("Angelina McMullen", "Coordinator", "")],
  "Plus 6 recruitment, 4 regulatory, 4 quality assurance, 2 phlebotomy and lab, and finance and contracting in house.",
  cols=3)
slide_panels("Facilities", "Both sites are equipped for Phase II to IV, so a protocol can open at either address.",
  "Philadelphia, PA — 9501 Roosevelt Blvd #208 · Pennington, NJ — 21 Route 31 North, Suite A8.",
  ("On-site equipment and lab",
   ["FibroScan, ultrasound, centrifuges, EKG, infusion pumps, incubator.",
    "Ambient, refrigerated 2–8°C and frozen −20°C / −80°C IP storage.",
    "24/7 temperature monitoring, sample processing lab, phlebotomy and nursing stations."], False),
  ("Sponsor and patient space",
   ["Private exam rooms and patient-only restrooms.",
    "Investigator and coordinator offices, dedicated monitor workstations.",
    "Sponsor conference and training room.",
    "Class-A medical building with 24-hour urgent care on campus, SEPTA accessible."], True))

# §5 track record
cover(4)
slide_stats("Executive snapshot", "Thirty years in the same two communities is what built the referral network.",
  "Two sites, more than 500 trials, and a referral network of 191 practices covering 1.5 million patients.",
  [("30", "years continuously operating since 1996"),
   ("500+", "clinical trials across diverse therapeutic areas"),
   ("191", "referring practices across PA and NJ"),
   ("200+", "years of combined research experience on team")],
  None, size=38)
slide_people("Investigators",
  "Eight investigators on staff, with neurology and geriatrics reached through our provider network.",
  "Combined 200+ years of clinical research experience.",
  [("Eugene Andruczyk, MBA, DO", "Investigator", "OB / GYN"),
   ("Lolita Vaughan, CRNP, CCRC", "Investigator", "Women's Health"),
   ("Joseph Heether, MD", "Investigator", "Cardiothoracic"),
   ("Brian Shaffer, MD", "Investigator", "Internal Medicine"),
   ("Vinod Rustgi, MD, MBA", "Investigator", "Hepatology / Liver"),
   ("Taher Modarressi, MD", "Investigator", "Endo / Lipidology"),
   ("Michael Tomeo, MD", "Investigator", "Dermatology"),
   ("Lawrence Leventhal, MD", "Investigator", "Rheumatology")],
  "Hal Ganzman, DO — sub-investigator, internal medicine.")
slide_table("Sponsor relationships", "Sponsors come back. Lilly has placed fifteen protocols.",
  "ICON already monitors protocols at our sites, including Lilly GZPO.",
  ["Sponsor", "Protocols", None, "Randomized", None],
  [["Eli Lilly", ("bar", 100, None, BLUE), "15", ("bar", 100, None), "80"],
   ["AbbVie", ("bar", 40, None, BLUE), "6", ("bar", 31, None), "25"],
   ["Johnson & Johnson", ("bar", 40, None, BLUE), "6", ("bar", 14, None), "11"],
   ["AstraZeneca", ("bar", 20, None, BLUE), "3", ("bar", 10, None), "8"]],
  "A sponsor who came back is the one reference a site cannot write for itself. Strategic partnerships with Lilly, J&J, AbbVie and AstraZeneca, including AstraZeneca's PIC program.",
  [3.2, 2.6, 1.0, 2.6, 1.0])
slide_stats("Track record", "Nine approved medicines came through trials we ran.",
  "Diabetes and obesity, migraine, and dermatology and rheumatology.",
  [("Wegovy · Ozempic · Mounjaro", "Diabetes and obesity"),
   ("Nurtec ODT · Ubrelvy · Qulipta", "Migraine"),
   ("Dupixent · Rinvoq · Olumiant", "Dermatology and rheumatology")],
  "Sponsors on the record: “CRP's commitment and professionalism is inspiring.” — IQVIA & Daiichi Sankyo. “Working with CRP's dedicated team has always been a great pleasure.” — Bayer.",
  size=17)

# §6 partnership
cover(5)
slide_panels("How we would operate", "Two locations. One contract, one system, one team.",
  "Philadelphia, PA and Pennington, NJ. Roughly 50 miles apart, drawing from separate catchments, run off a single system.",
  ("You contract once",
   ["A master agreement covers both sites — one budget process, one regulatory function.",
    "One EDC instance, so a monitor sees the same source structure at either address.",
    "One coordinator queue and one investigator roster across both."], False),
  ("You open at both addresses",
   ["Philadelphia and Pennington draw from separate catchments on one activation.",
    "24 protocols are enrolling today and neither site is at ceiling.",
    "Centralized operations, regulatory, recruitment, QA and finance."], True))

s = new("Put a number on us, and hold us to it.")
y = head(s, "What we will commit to", "Put a number on us, and hold us to it.")
items = [("1", "We commit to an enrollment number at feasibility, built from real-time mining of our whole patient database and broken down by therapeutic area — not a range, and not after the first month of screening."),
         ("2", "We report against it monthly, sent automatically to your study and CRO contacts."),
         ("3", "If we are tracking under by month two, you hear it from us first, with what we are changing.")]
cy = y + Inches(0.3)
for num, txt in items:
    circ = s.shapes.add_shape(K.MSO_SHAPE.OVAL, M, cy, Inches(0.42), Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
    circ.line.fill.background(); circ.shadow.inherit = False
    K._tf(circ, num, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, M + Inches(0.62), cy + Inches(0.02), CONTENT_W - Inches(0.62), Inches(0.7), txt, 14)
    cy += Inches(1.0)
ln = rect(s, M, cy + Inches(0.1), CONTENT_W, Pt(1.5), fill=NAVY, radius=False); ln.line.fill.background()
textbox(s, M, cy + Inches(0.32), CONTENT_W, Inches(0.4),
        "Aneesh Vaze, Managing Director · aneesh@phillyresearch.com · (215) 676-6696 · feasibility returned within 48 hours",
        12, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Our commitment", pageno())


# sponsor logo wall — the marks are a single rendered image (SVG will not embed
# in PowerPoint); the title, standfirst and quotes stay editable text.
_wall = DS / "assets/sponsors/_wall.png"
if _wall.exists():
    s = new("Large pharma and emerging biotech both place work here.")
    y = head(s, "Track record", "Large pharma and emerging biotech both place work here.",
             "A non-exhaustive selection.")
    s.shapes.add_picture(str(_wall), M, y + Inches(0.1), width=CONTENT_W)
    textbox(s, M, H - Inches(1.25), CONTENT_W, Inches(0.7),
            "\u201cCRP\u2019s commitment and professionalism is inspiring.\u201d  \u2014 IQVIA & Daiichi Sankyo\n"
            "\u201cWorking with CRP\u2019s dedicated team has always been a great pleasure.\u201d  \u2014 Bayer",
            11, color=GREY)
    footer(s, "Sponsors on the record", pageno())

prs.save(str(OUT))
print("wrote %s — %d editable slides, %.1f MB" % (OUT, PAGE[0], OUT.stat().st_size / 1e6))
